"""端到端冒烟测试：无需模型 Key / chromadb，验证 FastAPI + RAG + 对话链路。

用法（在 crop-agent-server 目录）：
  uv run python scripts/smoke_test.py
或（未使用 uv 时，先 pip install -r requirements.txt 后）：
  python scripts/smoke_test.py
"""
import os
import sys

# Windows 控制台默认 GBK，无法打印部分 Unicode 字符（如 ⚠）→ 容错输出
if hasattr(sys.stdout, "reconfigure"):
    sys.stdout.reconfigure(encoding="utf-8", errors="replace")
    sys.stderr.reconfigure(encoding="utf-8", errors="replace")

# 将 server 根目录加入 sys.path，保证 `import main` 与 `from config import ...` 可用
_SERVER_ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
sys.path.insert(0, _SERVER_ROOT)

# 必须在导入 main 之前设置环境（使用内存 SQLite，避免依赖文件系统）
os.environ["DB_URL"] = "sqlite:///:memory:"
os.environ["VECTOR_PERSIST_DIR"] = "./storage/smoke_chroma"
os.environ["REINDEX"] = "true"
# 冒烟测试保持确定性：关闭 LLM 查询改写（混合检索/重排不依赖网络）
os.environ["QUERY_REWRITE"] = "false"

# 知识库使用临时副本：管理类端点（删除/重建）会改文件，不能动真实 knowledge/
import shutil  # noqa: E402

_SMOKE_KB = os.path.abspath("./storage/smoke_knowledge")
if os.path.isdir(_SMOKE_KB):
    shutil.rmtree(_SMOKE_KB, ignore_errors=True)
shutil.copytree(os.path.abspath("../knowledge"), _SMOKE_KB)
os.environ["KNOWLEDGE_DIR"] = _SMOKE_KB

from fastapi.testclient import TestClient  # noqa: E402

import main  # noqa: E402


def check(name, cond, extra=""):
    print(f"[{'PASS' if cond else 'FAIL'}] {name} {extra}")
    return bool(cond)


ok = True
with TestClient(main.app) as client:
    r = client.get("/health")
    ok &= check("GET /health", r.status_code == 200, str(r.json()))

    r = client.post("/api/chat", json={"content": "水稻稻瘟病怎么防治？"})
    ok &= check(
        "POST /api/chat",
        r.status_code == 200 and r.json().get("code") == 200,
    )
    data = r.json()
    reply = data.get("data", {}).get("reply", "")
    sources = data.get("sources", [])
    print("  reply head:", reply[:80].replace("\n", " "))
    print("  sources:", [(s["chunk"], s["score"]) for s in sources])
    ok &= check("  检索命中来源", len(sources) >= 1)

    r = client.get("/api/knowledge/search", params={"q": "番茄早疫病"})
    ok &= check(
        "GET /api/knowledge/search",
        r.status_code == 200,
    )
    data = r.json().get("data", {})
    print("  results:", len(data.get("results", [])))

    # ---- 知识库界面优化：文档级合并 / 阈值 meta / 分类过滤 ----
    results = data.get("results", [])
    ok &= check(
        "  搜索结果按文档合并（含 chunks 数组）",
        all("chunks" in d and isinstance(d["chunks"], list) for d in results),
        repr([(d["doc_id"], len(d.get("chunks", []))) for d in results]),
    )
    meta = data.get("meta", {})
    ok &= check(
        "  返回有效阈值 meta",
        isinstance(meta.get("threshold"), float),
        repr(meta),
    )
    ok &= check(
        "  空查询返回空数组",
        client.get("/api/knowledge/search", params={"q": ""}).json()["data"]["results"]
        == [],
    )
    r = client.get("/api/knowledge/search", params={"q": "种植", "category": "diseases"})
    ok &= check(
        "  分类过滤生效",
        all(d["category"] == "diseases" for d in r.json()["data"]["results"]),
    )

    # ---- 混合检索（向量 + BM25 + RRF）：短查询 / 分数不变式 / rewrite 参数 ----
    r = client.get("/api/knowledge/search", params={"q": "稻瘟病", "rewrite": 0})
    short = r.json().get("data", {}).get("results", [])
    ok &= check(
        "  短查询'稻瘟病'混合检索有结果",
        r.status_code == 200 and len(short) >= 1,
        repr([(d["doc_id"], d["score"]) for d in short]),
    )
    all_scores = [d["score"] for d in results] + [d["score"] for d in short]
    ok &= check(
        "  展示分落在 [0,1]（加权融合分：缺失源记 0，单源命中不再被抬到双命中量级）",
        all(isinstance(s, float) and 0.0 <= s <= 1.0 for s in all_scores),
        repr(sorted(all_scores)),
    )
    ok &= check(
        "  结果附带后端判定的相关性等级",
        all(d.get("relevance") in ("high", "mid", "low", "none") for d in results),
        repr([(d["title"], d["score"], d.get("relevance")) for d in results[:3]]),
    )

    # 回归：逐字命中必须压过零星共字命中。曾因双命中片段的 BM25 分被静默丢弃，
    # 使向量分低到被阈值挡住、因而独占 BM25 分的《小麦》片段排到西瓜原文之前。
    melon_q = (
        "常见西瓜品种及特点 西瓜品种繁多，种植户应根据市场需求和当地气候条件"
        "选择合适的品种："
    )
    melon = client.get(
        "/api/knowledge/search", params={"q": melon_q, "rewrite": 0}
    ).json()["data"]["results"]
    if any("西瓜" in d["title"] for d in melon):
        ok &= check(
            "  西瓜原文查询：逐字命中排第一且判为高相关",
            "西瓜" in melon[0]["title"] and melon[0].get("relevance") == "high",
            repr([(d["title"], d["score"], d.get("relevance")) for d in melon[:3]]),
        )
    else:
        print("SKIP 西瓜排序回归断言（知识库无该文档，可能缺 python-docx）")
    r = client.get("/api/knowledge/search", params={"q": "番茄早疫病", "rewrite": 1})
    ok &= check(
        "  rewrite=1 参数不报错",
        r.status_code == 200,
        repr(r.json().get("message", "")),
    )

    r = client.get("/api/knowledge")
    docs = r.json()["data"]["docs"]
    ok &= check("GET /api/knowledge 含新字段", bool(docs) and all(
        k in docs[0] for k in ("status", "category", "file_size", "updated_at", "file_name")
    ), repr(docs[0]) if docs else "empty")
    ok &= check("  列表 meta 含阈值", "threshold" in r.json()["data"]["meta"])

    # 文档管理：chunks / content / reindex / delete / batch / chunk 删除
    if docs:
        doc_id = docs[0]["doc_id"]
        r = client.get(f"/api/knowledge/{doc_id}/chunks")
        chunk_list = r.json()["data"]["chunks"]
        ok &= check(
            "GET /{id}/chunks",
            r.status_code == 200 and len(chunk_list) >= 1 and "char_count" in chunk_list[0],
        )
        r = client.get(f"/api/knowledge/{doc_id}/content")
        ok &= check(
            "GET /{id}/content 原文预览",
            r.status_code == 200 and len(r.json()["data"]["content"]) > 0,
        )
        r = client.post(f"/api/knowledge/{doc_id}/reindex")
        ok &= check(
            "POST /{id}/reindex 单文档重建",
            r.json().get("code") == 200 and r.json()["data"]["chunk_count"] >= 1,
            repr(r.json()),
        )
        if len(chunk_list) > 1:
            cid = chunk_list[0]["chunk_id"]
            r = client.delete(f"/api/knowledge/chunks/{cid}")
            ok &= check(
                "DELETE /chunks/{id} 片段删除",
                r.json().get("code") == 200,
                repr(r.json()),
            )
            # 补回：重建索引恢复完整片段
            client.post(f"/api/knowledge/{doc_id}/reindex")
        elif chunk_list:
            cid = chunk_list[0]["chunk_id"]
            r = client.delete(f"/api/knowledge/chunks/{cid}")
            ok &= check(
                "  末个片段拒绝删除",
                r.json().get("code") == 400,
                repr(r.json()),
            )

        r = client.post("/api/knowledge/batch-reindex", json={"doc_ids": [doc_id]})
        ok &= check(
            "POST /batch-reindex",
            r.json().get("code") == 200 and r.json()["data"]["count"] == 1,
            repr(r.json()),
        )
        # 批量删除 + 单条删除（重建索引保证批量删除前文件存在）
        r = client.post("/api/knowledge/batch-delete", json={"doc_ids": [doc_id]})
        ok &= check(
            "POST /batch-delete",
            r.json().get("code") == 200 and r.json()["data"]["count"] == 1,
            repr(r.json()),
        )
        remain = [d["doc_id"] for d in client.get("/api/knowledge").json()["data"]["docs"]]
        ok &= check("  批量删除后列表不含", doc_id not in remain)
        r = client.delete(f"/api/knowledge/{doc_id}")
        ok &= check(
            "  已删除文档再删返回 404",
            r.json().get("code") == 404,
            repr(r.json()),
        )

        # 上传链路：新建临时文档 → 上传 → 列表出现 → 清理
        tmp_path = os.path.join(_SMOKE_KB, "冒烟测试临时文档.md")
        with open(tmp_path, "w", encoding="utf-8") as f:
            f.write("# 冒烟测试临时文档\n\n这是一份用于验证上传与单文档索引链路的文档。\n" * 5)
        with open(tmp_path, "rb") as f:
            r = client.post(
                "/api/knowledge/upload",
                files={"file": ("冒烟测试临时文档.md", f, "text/markdown")},
            )
        up_data = r.json()
        ok &= check(
            "POST /api/knowledge/upload",
            up_data.get("code") == 200 and up_data["data"].get("doc_id"),
            repr(up_data),
        )
        up_id = up_data["data"]["doc_id"]
        now_ids = [d["doc_id"] for d in client.get("/api/knowledge").json()["data"]["docs"]]
        ok &= check("  上传后列表包含新文档", up_id in now_ids)
        r = client.delete(f"/api/knowledge/{up_id}")
        ok &= check(
            "  上传文档可删除",
            r.json().get("code") == 200,
            repr(r.json()),
        )
        if os.path.isfile(tmp_path):
            os.remove(tmp_path)

    # ---- 文档管理：修改分类 + 与检索测试同步（需求3） ----
    cat_path = os.path.join(_SMOKE_KB, "冒烟测试草莓.md")
    with open(cat_path, "w", encoding="utf-8") as f:
        f.write("# 冒烟测试草莓\n\n草莓种植要点关键词：草莓喜光、适宜温度 15-25℃，需注意白粉病防治。\n" * 3)
    with open(cat_path, "rb") as f:
        r = client.post(
            "/api/knowledge/upload",
            files={"file": ("冒烟测试草莓.md", f, "text/markdown")},
        )
    cat_up = r.json()
    ok &= check(
        "  上传分类测试文档",
        cat_up.get("code") == 200 and cat_up["data"].get("doc_id"),
        repr(cat_up),
    )
    cat_id = cat_up["data"]["doc_id"]

    r = client.patch(f"/api/knowledge/{cat_id}", json={"category": "crops"})
    ok &= check(
        "PATCH /{id} 修改分类",
        r.json().get("code") == 200 and r.json()["data"]["category"] == "crops",
        repr(r.json()),
    )

    def _find_doc(q, category=""):
        params = {"q": q}
        if category:
            params["category"] = category
        return client.get("/api/knowledge/search", params=params).json()["data"]["results"]

    ok &= check(
        "  改分类后无过滤检索命中",
        any(d["doc_id"] == cat_id for d in _find_doc("草莓种植要点")),
    )
    ok &= check(
        "  改分类后按新分类 crops 命中",
        any(d["doc_id"] == cat_id and d["category"] == "crops" for d in _find_doc("草莓种植要点", "crops")),
    )
    ok &= check(
        "  改分类后按旧分类 diseases 不命中（向量元数据已同步）",
        all(d["doc_id"] != cat_id for d in _find_doc("草莓种植要点", "diseases")),
    )
    list_doc = next(
        (d for d in client.get("/api/knowledge").json()["data"]["docs"] if d["doc_id"] == cat_id),
        None,
    )
    ok &= check(
        "  列表分类已更新为 crops",
        bool(list_doc) and list_doc["category"] == "crops",
        repr(list_doc),
    )

    r = client.patch(f"/api/knowledge/{cat_id}", json={"category": "   "})
    ok &= check("  空分类被拒绝(400)", r.json().get("code") == 400, repr(r.json()))
    r = client.patch(f"/api/knowledge/{cat_id}", json={"category": "x" * 65})
    ok &= check("  超长分类被拒绝(400)", r.json().get("code") == 400, repr(r.json()))
    r = client.patch("/api/knowledge/NONEXIST", json={"category": "crops"})
    ok &= check("  不存在文档改分类返回 404", r.json().get("code") == 404, repr(r.json()))

    r = client.delete(f"/api/knowledge/{cat_id}")
    ok &= check("  分类测试文档可删除", r.json().get("code") == 200, repr(r.json()))
    ok &= check(
        "  删除后检索测试不再命中（同步生效）",
        all(d["doc_id"] != cat_id for d in _find_doc("草莓种植要点")),
    )
    if os.path.isfile(cat_path):
        os.remove(cat_path)

    # ---- Word/PPT 支持：解析 / 分段检索 / 预览 / 下载 ----
    try:
        import docx  # noqa: F401
        import pptx  # noqa: F401
    except ImportError as exc:
        print(f"SKIP docx/pptx checks（缺解析库）：{exc}")
        has_office = False
    else:
        has_office = True

    if has_office:
        from docx import Document
        from pptx import Presentation

        # 生成临时 .docx / .pptx 夹具（含唯一关键词）
        docx_path = os.path.join(_SMOKE_KB, "冒烟测试Word.docx")
        d = Document()
        d.add_paragraph("冒烟测试Word文档关键词")
        d.add_paragraph("水稻稻瘟病防治要点（来自 Word 文档）")
        d.save(docx_path)

        prs = Presentation()
        layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]
        slide = prs.slides.add_slide(layout)
        tf = slide.shapes.add_textbox(0, 0, 9144000, 914400).text_frame
        tf.text = "冒烟测试PPT关键词"
        tf.add_paragraph().text = "番茄早疫病的识别与防治（来自 PPT 文档）"
        pptx_path = os.path.join(_SMOKE_KB, "冒烟测试PPT.pptx")
        prs.save(pptx_path)

        for fname, keyword, fmt in [
            ("冒烟测试Word.docx", "冒烟测试Word文档关键词", "docx"),
            ("冒烟测试PPT.pptx", "冒烟测试PPT关键词", "pptx"),
        ]:
            path = os.path.join(_SMOKE_KB, fname)
            with open(path, "rb") as f:
                r = client.post(
                    "/api/knowledge/upload",
                    files={"file": (fname, f, "application/octet-stream")},
                )
            up = r.json()
            ok &= check(
                f"上传 {fmt} 成功",
                up.get("code") == 200 and up["data"].get("doc_id"),
                repr(up),
            )
            fid = up["data"]["doc_id"]
            r = client.get("/api/knowledge/search", params={"q": keyword})
            hits = r.json().get("data", {}).get("results", [])
            hit = next((x for x in hits if x["doc_id"] == fid), None)
            ok &= check(
                f"{fmt} 分段检索命中（含 chunks）",
                hit is not None and len(hit.get("chunks", [])) >= 1,
                repr(hit),
            )
            r = client.get(f"/api/knowledge/{fid}/content")
            body = r.json().get("data", {})
            ok &= check(
                f"{fmt} 内容预览（format + 提取文本）",
                r.json().get("code") == 200
                and body.get("format") == fmt
                and keyword in body.get("content", ""),
                repr(body)[:200],
            )
            r = client.get(f"/api/knowledge/{fid}/download")
            cd = r.headers.get("content-disposition", "")
            ok &= check(
                f"{fmt} 下载原始文件",
                r.status_code == 200 and "attachment" in cd and "filename" in cd,
                cd[:120],
            )
            r = client.delete(f"/api/knowledge/{fid}")
            ok &= check(f"{fmt} 上传后可删除", r.json().get("code") == 200, repr(r.json()))
            if os.path.isfile(path):
                os.remove(path)

    # 旧格式 .doc 拒绝并提示转换
    r = client.post(
        "/api/knowledge/upload",
        files={"file": ("旧格式.doc", b"\xd0\xcf\xd2\xc2 not a real doc", "application/msword")},
    )
    msg = r.json().get("message", "")
    ok &= check(".doc 旧格式被拒绝并提示转换", r.json().get("code") == 400 and ".docx" in msg, msg)

    r = client.get("/api/sessions")
    sessions = r.json()["data"]["sessions"]
    ok &= check(
        "GET /api/sessions",
        r.status_code == 200 and len(sessions) >= 1,
    )
    ok &= check(
        "  会话标题取首轮提问",
        bool(sessions) and sessions[0]["title"].startswith("水稻稻瘟病"),
        repr(sessions[0]["title"] if sessions else None),
    )

    sid = sessions[0]["session_id"] if sessions else ""

    # ---- 会话管理新能力：重命名 / 置顶 / 批量删除 ----
    r = client.patch(f"/api/sessions/{sid}", json={"title": "重命名后的对话"})
    ok &= check(
        "PATCH /api/sessions/{id} 重命名",
        r.status_code == 200 and r.json()["data"]["title"] == "重命名后的对话",
        repr(r.json()),
    )

    r = client.patch(f"/api/sessions/{sid}", json={"pinned": True})
    ok &= check(
        "PATCH /api/sessions/{id} 置顶",
        r.status_code == 200 and r.json()["data"]["pinned"] is True,
    )
    r = client.get("/api/sessions")
    top = r.json()["data"]["sessions"][0]
    ok &= check(
        "  置顶会话排最前",
        top["session_id"] == sid and top["pinned"] is True,
        repr(top),
    )
    r = client.patch(f"/api/sessions/{sid}", json={"pinned": False})
    ok &= check("  取消置顶", r.json()["data"]["pinned"] is False)

    r = client.patch(f"/api/sessions/{sid}", json={"title": "   "})
    ok &= check(
        "  空标题被拒绝",
        r.json().get("code") == 400,
        repr(r.json()),
    )

    # 批量删除：新建 2 个会话后一次性删除
    batch_ids = []
    for _ in range(2):
        r = client.post("/api/sessions")
        batch_ids.append(r.json()["data"]["session_id"])
    r = client.post("/api/sessions/batch-delete", json={"session_ids": batch_ids})
    ok &= check(
        "POST /api/sessions/batch-delete",
        r.status_code == 200 and r.json()["data"]["count"] == 2,
        repr(r.json()),
    )
    r = client.get("/api/sessions")
    remain = [s["session_id"] for s in r.json()["data"]["sessions"]]
    ok &= check(
        "  批量删除后列表不含",
        all(i not in remain for i in batch_ids),
        repr(remain),
    )

    r = client.delete(f"/api/sessions/{sid}")
    ok &= check(
        "DELETE /api/sessions/{id}",
        r.status_code == 200 and r.json()["data"]["deleted"] == sid,
    )
    r = client.get("/api/chat/history", params={"session_id": sid})
    msgs = r.json()["data"]["messages"] if r.status_code == 200 else None
    ok &= check(
        "  删除后消息清空",
        msgs == [],
        repr(msgs),
    )

print("\nSMOKE", "PASS" if ok else "FAIL")
sys.exit(0 if ok else 1)
